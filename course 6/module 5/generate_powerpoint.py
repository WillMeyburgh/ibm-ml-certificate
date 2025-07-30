from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()

    # Slide 1: Cover Page
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Data Findings Report: Course Recommendation System"
    subtitle.text = "IBM Machine Learning Professional Certificate Capstone Project\nJuly 29, 2025"

    # Slide 2: Executive Summary
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Executive Summary"
    executive_summary_text = """
    Initial EDA revealed key insights into course popularity, genre distribution, and user enrollment.
    Feature engineering (Bag of Words) effectively transformed text data.
    Both content-based and collaborative filtering showed promising results.
    Hybrid approaches and dimensionality reduction (PCA, NMF) enhanced accuracy and scalability.
    """
    content.text = executive_summary_text.strip()
    content.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    # Slide 3: Table of Contents
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Table of Contents"
    toc_items = [
        "1. Introduction",
        "2. Exploratory Data Analysis",
        "3. Feature Engineering",
        "4. Recommendation System Implementations",
        "5. Evaluation",
        "6. Innovative Insights",
        "7. Conclusion"
    ]
    for item in toc_items:
        p = content.text_frame.add_paragraph()
        p.text = item
        p.level = 0

    # Slide 4: Introduction - Project Goal & Motivation
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Introduction: Project Goal & Motivation"
    intro_text = """
    Project Goal: To develop a robust course recommendation system for an online learning platform.
    Motivation: Enhance user engagement and learning paths by providing personalized course suggestions.
    """
    content.text = intro_text.strip()
    content.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    # Slide 5: Introduction - Methodology Overview & Target Audience
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Introduction: Methodology & Target Audience"
    methodology_text = """
    Methodology Overview: This report details the process of building a recommendation engine, covering data analysis, feature engineering, and the implementation and evaluation of various recommendation algorithms, including content-based filtering, collaborative filtering (KNN, NMF, and neural network embeddings), and hybrid approaches.
    Target Audience: This report is intended for stakeholders, data scientists, and product managers interested in understanding the development and performance of the course recommendation system.
    """
    content.text = methodology_text.strip()
    content.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    # Slide 6: Introduction - Report Structure
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Introduction: Report Structure"
    report_structure_items = [
        "Introduction: Overview of project, goals, and methodology.",
        "Exploratory Data Analysis: Insights from initial data exploration.",
        "Feature Engineering: Preparing data for model training.",
        "Recommendation System Implementations: Algorithms used and results.",
        "Evaluation: Performance comparison of approaches.",
        "Conclusion: Project outcomes and future work."
    ]
    for item in report_structure_items:
        p = content.text_frame.add_paragraph()
        p.text = item
        p.level = 0

    # Slide 7: Exploratory Data Analysis - Course Titles and Genres
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Exploratory Data Analysis: Course Titles & Genres"
    content.text = "Keyword Analysis: Word cloud reveals focus on IT skills (Python, data science, machine learning, big data, AI, TensorFlow, cloud, containers)."
    content.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    p = content.text_frame.add_paragraph()
    p.text = "Genre Distribution: BackendDev, MachineLearning, Database, DataAnalysis are most frequent categories."

    # Slide 8: Exploratory Data Analysis - Course Titles and Genres (Images)
    slide_layout = prs.slide_layouts[5] # Title Only layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Course Titles & Genres Visualizations"

    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(4.5)
    height = Inches(5)
    try:
        pic = slide.shapes.add_picture("/home/annekin/workspace/ibm-ml-certificate/course 6/module 5/assets/eda_wordcloud.png", left, top, width, height)
    except FileNotFoundError:
        print("Warning: eda_wordcloud.png not found. Skipping image.")

    left = Inches(5)
    top = Inches(1.5)
    width = Inches(4.5)
    height = Inches(5)
    try:
        pic = slide.shapes.add_picture("/home/annekin/workspace/ibm-ml-certificate/course 6/module 5/assets/eda_genrecounts.png", left, top, width, height)
    except FileNotFoundError:
        print("Warning: eda_genrecounts.png not found. Skipping image.")

    # Slide 8: Exploratory Data Analysis - Course Enrollments
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Exploratory Data Analysis: Course Enrollments"
    content.text = "User Enrollment Statistics: 233,306 total enrollments from 33,901 unique users. Avg. 7 courses/user, max 61."
    content.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    p = content.text_frame.add_paragraph()
    p.text = "Top Courses: Top 20 courses account for >63% of enrollments. 'Introduction to Python' is most popular (>15,000 enrollments)."

    # Slide 9: Exploratory Data Analysis - Course Enrollments (Image)
    slide_layout = prs.slide_layouts[5] # Title Only layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Course Enrollments Visualization"

    left = Inches(2)
    top = Inches(1.5)
    width = Inches(6)
    height = Inches(5)
    try:
        pic = slide.shapes.add_picture("/home/annekin/workspace/ibm-ml-certificate/course 6/module 5/assets/eda_hist.png", left, top, width, height)
    except FileNotFoundError:
        print("Warning: eda_hist.png not found. Skipping image.")

    # Slide 9: Feature Engineering (Bag of Words)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Feature Engineering: Bag of Words (BoW)"
    content.text = "To prepare course text data for machine learning, Bag of Words (BoW) features were extracted from course titles and descriptions."
    content.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    
    p = content.text_frame.add_paragraph()
    p.text = "The process involved:"
    p.level = 0
    
    p = content.text_frame.add_paragraph()
    p.text = "Tokenization: Splitting the text into individual words (tokens)."
    p.level = 1
    
    p = content.text_frame.add_paragraph()
    p.text = "Stop Word Removal: Filtering out common, non-informative words."
    p.level = 1
    
    p = content.text_frame.add_paragraph()
    p.text = "Part-of-Speech (POS) Tagging: Identifying and keeping only nouns."
    p.level = 1

    p = content.text_frame.add_paragraph()
    p.text = "Result: Each course represented by a vector of token counts."
    p.level = 0

    # Slide 10: Course Similarity
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Course Similarity: Cosine Similarity"
    content.text = "Cosine Similarity: BoW feature vectors used to measure likeness between courses."
    content.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    p = content.text_frame.add_paragraph()
    p.text = "Example: Finding courses similar to 'Machine Learning with Python' based on high cosine similarity scores."
    p = content.text_frame.add_paragraph()
    p.text = "This forms the basis for a content-based recommender system."

    # Slide 11: Content-Based Recommender System (User Profile & Course Genres) - Methodology
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Content-Based RS: User Profile & Genres (Methodology)"
    content.text = "User profiles are generated by creating a weighted genre vector. This vector is calculated by multiplying a user's course ratings with the genre vectors of the courses they have rated."
    content.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    p = content.text_frame.add_paragraph()
    p.text = "Recommendation scores are computed by taking the dot product of a user's profile vector and the genre vector of a course they have not yet taken."

    # Slide 12: Content-Based Recommender System (User Profile & Course Genres) - Results & Visualization
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Content-Based RS: User Profile & Genres (Results)"
    content.text = "The system successfully generates personalized course recommendations. For example, a user with a strong interest in Python and Machine Learning receives recommendations for courses like Python 101 and Machine Learning with R."
    content.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    # Slide 15: Content-Based Recommender System (User Profile & Course Genres) - Visualization
    slide_layout = prs.slide_layouts[5] # Title Only layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "User Profile & Course Genres Visualization"

    left = Inches(2)
    top = Inches(1.5)
    width = Inches(6)
    height = Inches(5)
    try:
        pic = slide.shapes.add_picture("/home/annekin/workspace/ibm-ml-certificate/course 6/module 5/assets/user_profile_bar.png", left, top, width, height)
    except FileNotFoundError:
        print("Warning: user_profile_bar.png not found. Skipping image.")

    # Slide 13: Content-Based Recommender System (Course Similarity) - Methodology
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Content-Based RS: Course Similarity (Methodology)"
    content.text = "A course-to-course similarity matrix is computed using the Bag of Words (BoW) features of each course. The similarity score ranges from 0 to 1, where 1 indicates a perfect match."
    content.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    p = content.text_frame.add_paragraph()
    p.text = "To recommend courses, the system takes a user's list of enrolled courses and a similarity threshold (e.g., > 0.6) as input. It iterates through each enrolled course and finds other courses with a similarity score above the threshold. The system aggregates all unique, similar courses that the user has not already taken and sorts them by their similarity score to generate a ranked list of recommendations."

    # Slide 14: Content-Based Recommender System (Course Similarity) - Results & Visualization
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Content-Based RS: Course Similarity (Results)"
    content.text = "The system effectively recommends new courses based on a user's enrollment history. For instance, if a user is enrolled in machine learning courses, the recommender suggests other relevant courses in data science, deep learning, and Python."
    content.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    p = content.text_frame.add_paragraph()
    p.text = "Example: Given a set of enrolled courses, the system identifies and ranks 20 new courses with a similarity greater than 0.6."

    # Slide 18: Content-Based Recommender System (Course Similarity) - Visualization
    slide_layout = prs.slide_layouts[5] # Title Only layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Course Similarity Visualization"

    left = Inches(2)
    top = Inches(1.5)
    width = Inches(6)
    height = Inches(5)
    try:
        pic = slide.shapes.add_picture("/home/annekin/workspace/ibm-ml-certificate/course 6/module 5/assets/course_similarity_heatmap.png", left, top, width, height)
    except FileNotFoundError:
        print("Warning: course_similarity_heatmap.png not found. Skipping image.")

    # Slide 15: Content-Based Recommender System (User Profile Clustering) - Methodology
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Content-Based RS: User Profile Clustering (Methodology)"
    content.text = "This approach groups users with similar interests into clusters and then recommends popular courses from within those clusters. Two methods were tested:"
    content.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    
    p = content.text_frame.add_paragraph()
    p.text = "1. Clustering on Original Features: K-Means on 14 standardized user-profile features. Optimal k=10 (elbow method)."
    p.level = 1
    
    p = content.text_frame.add_paragraph()
    p.text = "2. Clustering on PCA-Reduced Features: PCA reduced dimensionality (9 components explain >90% variance). K-Means on 9 components, optimal k=20."
    p.level = 1

    p = content.text_frame.add_paragraph()
    p.text = "For a given user, the system identifies their cluster and recommends the most frequently enrolled courses by other users in the same cluster, excluding courses the user has already taken."
    p.level = 0

    # Slide 16: Content-Based Recommender System (User Profile Clustering) - Results & Visualizations
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Content-Based RS: User Profile Clustering (Results)"
    content.text = "Both methods successfully created meaningful user clusters, or 'learning communities,' based on shared interests."
    content.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    p = content.text_frame.add_paragraph()
    p.text = "When tested on a sample user, both approaches generated relevant course recommendations. The PCA-based method produced comparable recommendations while being more computationally efficient."

    # Slide 21: Content-Based Recommender System (User Profile Clustering) - Visualizations
    slide_layout = prs.slide_layouts[5] # Title Only layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "User Profile Clustering Visualizations"

    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(4.5)
    height = Inches(5)
    try:
        pic = slide.shapes.add_picture("/home/annekin/workspace/ibm-ml-certificate/course 6/module 5/assets/cluster_elbow.png", left, top, width, height)
    except FileNotFoundError:
        print("Warning: cluster_elbow.png not found. Skipping image.")

    left = Inches(5)
    top = Inches(1.5)
    width = Inches(4.5)
    height = Inches(5)
    try:
        pic = slide.shapes.add_picture("/home/annekin/workspace/ibm-ml-certificate/course 6/module 5/assets/cluster_covariance.png", left, top, width, height)
    except FileNotFoundError:
        print("Warning: cluster_covariance.png not found. Skipping image.")

    # Slide 17: KNN-Based Collaborative Filtering
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "KNN-Based Collaborative Filtering"
    content.text = "Overview: Collaborative filtering (CF) is a widely used recommendation approach, categorized into user-based (finding similar users) and item-based (finding similar items) methods."
    content.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    p = content.text_frame.add_paragraph()
    p.text = "KNN Approach: User-based (calculates similarity between user rating vectors to identify k nearest neighbors) or Item-based (calculates similarity between item rating vectors to find similar items)."
    p = content.text_frame.add_paragraph()
    p.text = "Example: A predicted rating of 2.77 for a course with a true rating of 3.0 resulted in an RMSE of 0.22, indicating high prediction accuracy."
    p = content.text_frame.add_paragraph()
    p.text = "Conclusion: KNN-based CF is effective but can be memory-intensive due to large similarity matrices. Future work may explore less memory-demanding CF approaches."

    # Slide 18: NMF-Based Collaborative Filtering - Overview & Methodology
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "NMF-Based Collaborative Filtering (Overview & Methodology)"
    content.text = "Overview: Non-negative Matrix Factorization (NMF) is a dimensionality reduction technique used to address the scalability issues of memory-based collaborative filtering methods like KNN."
    content.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    p = content.text_frame.add_paragraph()
    p.text = "Methodology: NMF decomposes a large, sparse user-item interaction matrix into two smaller, dense matrices: A user feature matrix (U) and an item feature matrix (I)."
    p = content.text_frame.add_paragraph()
    p.text = "The product of a user's latent feature vector from U and an item's latent feature vector from I provides an estimation of the original rating. The values in U and I are optimized by minimizing a cost function (e.g., squared difference between actual and estimated ratings) using optimization algorithms like Stochastic Gradient Descent (SGD)."

    # Slide 19: NMF-Based Collaborative Filtering - Results
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "NMF-Based Collaborative Filtering (Results)"
    content.text = "The NMF model achieved a Root Mean Squared Error (RMSE) of 1.3078 on the test set, indicating its predictive accuracy in estimating user ratings."
    content.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    # Slide 20: Neural Network Embedding-Based Collaborative Filtering - Overview & Results
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Neural Network Embedding-Based CF (Overview & Results)"
    content.text = "Overview: Neural networks can be effectively used to extract latent user and item features, similar to NMF, for collaborative filtering. This approach allows for rating prediction without explicitly pre-building feature vectors."
    content.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    p = content.text_frame.add_paragraph()
    p.text = "Results: The model achieved a root_mean_squared_error of approximately 0.0894 on the training set and 0.1218 on the validation set after 10 epochs, demonstrating its ability to predict course ratings."
    p = content.text_frame.add_paragraph()
    p.text = "The weights learned by the user_embedding_layer and item_embedding_layer within the trained neural network represent the extracted latent features for users and items (e.g., 16-dimensional latent feature vector for each user/item)."

    # Slide 21: Neural Network Embedding-Based Collaborative Filtering - Visualization
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    
    title.text = "Neural Network Embedding-Based CF (Visualization)"
    # Slide 26: Neural Network Embedding-Based Collaborative Filtering - Visualization
    slide_layout = prs.slide_layouts[5] # Title Only layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Neural Network Embedding-Based CF (Visualization)"

    left = Inches(2)
    top = Inches(1.5)
    width = Inches(6)
    height = Inches(5)
    try:
        pic = slide.shapes.add_picture("/home/annekin/workspace/ibm-ml-certificate/course 6/module 5/assets/nn_loss.png", left, top, width, height)
    except FileNotFoundError:
        print("Warning: nn_loss.png not found. Skipping image.")

    # Slide 22: Collaborative Filtering Algorithms Evaluation - KNN & NMF
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Collaborative Filtering Algorithms Evaluation (KNN & NMF)"
    content.text = "KNN-based Collaborative Filtering:"
    content.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    p = content.text_frame.add_paragraph()
    p.text = "Metric: Root Mean Squared Error (RMSE)."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "Result Example: Predicted rating of 2.77 (true 3.0) resulted in RMSE of 0.22 (high accuracy). Note: Effective but can be memory-intensive."
    p.level = 1

    p = content.text_frame.add_paragraph()
    p.text = "NMF-based Collaborative Filtering:"
    p.level = 0
    p = content.text_frame.add_paragraph()
    p.text = "Metric: Root Mean Squared Error (RMSE)."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "Result: NMF model achieved an RMSE of 1.3078 on the test set."
    p.level = 1

    # Slide 23: Collaborative Filtering Algorithms Evaluation - Neural Network Embedding
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Collaborative Filtering Algorithms Evaluation (NN Embedding)"
    content.text = "Neural Network Embedding-based Collaborative Filtering:"
    content.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    p = content.text_frame.add_paragraph()
    p.text = "Metric: Root Mean Squared Error (RMSE)."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "Results: The model achieved an RMSE of approximately 0.0894 on the training set and 0.1218 on the validation set after 10 epochs."
    p.level = 1

    # Slide 29: Collaborative Filtering Algorithms Evaluation - Neural Network Embedding (Visualization)
    slide_layout = prs.slide_layouts[5] # Title Only layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Neural Network Embedding-Based CF Evaluation Visualization"

    left = Inches(2)
    top = Inches(1.5)
    width = Inches(6)
    height = Inches(5)
    try:
        pic = slide.shapes.add_picture("/home/annekin/workspace/ibm-ml-certificate/course 6/module 5/assets/nn_loss.png", left, top, width, height)
    except FileNotFoundError:
        print("Warning: nn_loss.png not found. Skipping image.")

    # Slide 24: Innovative Insights - Hybrid Approach & Interpretability
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Innovative Insights: Hybrid Approach & Interpretability"
    content.text = "Hybrid Recommendation Approach: The project demonstrates the effectiveness of combining content-based and collaborative filtering methods. This hybrid approach leverages the strengths of both, providing robust recommendations even for new users or items (cold-start problem) by utilizing content features, while also capturing complex user-item interactions through collaborative filtering."
    content.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    p = content.text_frame.add_paragraph()
    p.text = "Interpretability of Latent Features: Through techniques like NMF and Neural Network Embeddings, the project successfully extracts latent features for users and courses. These features, while abstract, provide a lower-dimensional representation of underlying preferences and characteristics, offering a more interpretable understanding of why certain recommendations are made."

    # Slide 25: Innovative Insights - Scalability & Actionable Insights
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Innovative Insights: Scalability & Actionable Insights"
    content.text = "Scalability through Dimensionality Reduction: The application of PCA for user clustering and NMF for collaborative filtering highlights the importance of dimensionality reduction in handling large datasets. These techniques not only improve computational efficiency but also help in identifying the most significant patterns in the data, leading to more effective and scalable recommendation systems."
    p = content.text_frame.add_paragraph()
    p.text = "Actionable Insights for Course Development: The analysis of course genres and popular keywords provides direct, actionable insights for course developers. Understanding the most sought-after topics (e.g., Python, Machine Learning, AI) can guide the creation of new courses or the enhancement of existing ones to meet market demand."

    # Slide 26: Conclusion
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Conclusion"
    conclusion_items = [
        "Successfully developed a robust course recommendation system.",
        "Explored various techniques: content-based, collaborative filtering (KNN, NMF, NN embeddings), and hybrid approaches.",
        "Demonstrated effectiveness of feature engineering and dimensionality reduction.",
        "Future work: Explore less memory-demanding CF, real-time recommendations, and A/B testing."
    ]
    for item in conclusion_items:
        p = content.text_frame.add_paragraph()
        p.text = item
        p.level = 0

    prs.save("/home/annekin/workspace/ibm-ml-certificate/course 6/module 5/Data_Findings_Report.pptx")

if __name__ == "__main__":
    create_presentation()
